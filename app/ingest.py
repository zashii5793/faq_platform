"""ナレッジ取り込みパイプライン。

責務:
  1. ファイルパース（CSV / Markdown / テキスト / JSON / PDF / Excel / PowerPoint）
  2. チャンク化
  3. PII・機密マーカー検出
  4. 推奨アクション判定（ok / warn / danger）
  5. FAQマスターへの書き込み

注意:
  - 画像OCR は ROADMAP Phase 2 で対応予定
"""
from __future__ import annotations

import hashlib
import io
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

import pandas as pd

from .masking import build_rules, mask

Recommendation = Literal["ok", "warn", "danger"]

# 機密マーカー
CONFIDENTIAL_MARKERS = [
    "社外秘", "極秘", "[CONFIDENTIAL]", "Confidential",
    "禁転載", "取扱注意", "marked confidential",
]

# 個人氏名らしさ（簡易ヒューリスティック）
# 苗字 + 空白 + 名前 のパターンを要求（false positive 抑制）
# 例: "山田 太郎", "佐藤 花子", "ヤマダ タロウ"
# CSV等で空白なしの場合は別シグナル（mynumber等）に委ねる
NAME_HEURISTIC = re.compile(
    r"[一-鿿]{1,4}[ 　]+[一-鿿]{1,4}"
    r"|[ァ-ヴー]{2,5}[ 　]+[ァ-ヴー]{2,5}"
)


@dataclass
class Chunk:
    chunk_id: str
    source: str
    text: str


@dataclass
class ChunkFindings:
    """個別チャンクの検出結果。"""
    pii_counts: dict[str, int] = field(default_factory=dict)
    confidential_markers: list[str] = field(default_factory=list)
    name_candidates: int = 0
    recommendation: Recommendation = "ok"
    reason: str = ""


@dataclass
class FileFindings:
    pii_counts: dict[str, int] = field(default_factory=dict)
    confidential_markers: list[str] = field(default_factory=list)
    name_candidates: int = 0


@dataclass
class FileAnalysis:
    filename: str
    sha256: str
    size_bytes: int
    format: str
    chunks: list[Chunk]
    findings: FileFindings
    recommendation: Recommendation
    reason: str
    chunk_findings: list[ChunkFindings] = field(default_factory=list)

    @property
    def n_chunks(self) -> int:
        return len(self.chunks)


# =====================================================
# パース
# =====================================================
def _detect_format(filename: str) -> str:
    ext = Path(filename).suffix.lower().lstrip(".")
    return {
        "md": "markdown",
        "markdown": "markdown",
        "txt": "text",
        "csv": "csv",
        "json": "json",
        "pdf": "pdf",
        "xlsx": "xlsx",
        "xls": "xlsx",
        "pptx": "pptx",
        "ppt": "pptx",
    }.get(ext, "unsupported")


def _split_text(text: str, max_chars: int = 350) -> list[str]:
    """テキストをチャンクに分割。

    分割ルール:
      1. Markdown 見出し (## / ### で始まる行) は必ず新チャンクの境界
      2. それ以外は段落（空行2つ）単位で結合し、max_chars を超えたら分割
      3. 1つのチャンクが max_chars * 1.5 を超える場合は強制改行
      4. 空白のみのチャンクは破棄
    """
    lines = text.split("\n")
    sections: list[str] = []
    buf: list[str] = []
    for line in lines:
        if line.startswith(("## ", "### ")) and buf:
            sections.append("\n".join(buf).strip())
            buf = [line]
        else:
            buf.append(line)
    if buf:
        sections.append("\n".join(buf).strip())
    sections = [s for s in sections if s]

    chunks: list[str] = []
    for sec in sections:
        if len(sec) <= max_chars * 1.5:
            chunks.append(sec)
            continue
        # セクションが長すぎる場合は段落単位で分割
        paragraphs = [p.strip() for p in sec.split("\n\n") if p.strip()]
        cur = ""
        for p in paragraphs:
            if len(cur) + len(p) + 2 <= max_chars:
                cur = (cur + "\n\n" + p).strip()
            else:
                if cur:
                    chunks.append(cur)
                cur = p
        if cur:
            chunks.append(cur)
    # 最終フィルタ: 空・空白のみのチャンクは破棄
    return [c for c in chunks if c.strip()]


def _safe_filename(filename: str) -> str:
    """アップロードファイル名を安全化:
    - ディレクトリトラバーサル（../, ./）を防ぐためベース名のみ取得
    - スラッシュ・バックスラッシュを除去
    - 制御文字を除去
    - 空文字なら 'unnamed.txt'
    """
    import os
    # パス区切りを除去（path traversal対策）
    safe = os.path.basename(filename.replace("\\", "/").replace("\0", ""))
    # 制御文字除去
    safe = "".join(c for c in safe if ord(c) >= 32)
    safe = safe.strip()
    return safe or "unnamed.txt"


def _parse_text(filename: str, content: str) -> list[Chunk]:
    return [
        Chunk(chunk_id=f"{filename}#{i}", source=filename, text=t)
        for i, t in enumerate(_split_text(content))
    ]


def _parse_csv(filename: str, content: bytes) -> list[Chunk]:
    """CSV: 各行を 1 チャンクに変換（ヘッダーつき "key: value" 形式）。"""
    df = pd.read_csv(io.BytesIO(content), encoding_errors="replace")
    chunks: list[Chunk] = []
    for i, row in df.iterrows():
        parts = [f"{col}: {row[col]}" for col in df.columns if pd.notna(row[col])]
        text = " / ".join(parts)
        chunks.append(Chunk(chunk_id=f"{filename}#row{i+1}", source=filename, text=text))
    return chunks


def _parse_pdf(filename: str, content: bytes) -> list[Chunk]:
    """PDF: ページ単位でチャンク化（chunk_id にページ番号を含める）。

    壊れた PDF・パースできない PDF は例外を握って空リストを返す。
    （呼び出し側の analyze() で「テキスト抽出できませんでした」として警告表示する）
    """
    try:
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError
    except ImportError:
        return []

    try:
        reader = PdfReader(io.BytesIO(content))
    except (PdfReadError, ValueError, OSError, Exception):
        # 壊れた PDF / PDF でない / 暗号化済み 等
        return []

    chunks: list[Chunk] = []
    for i, page in enumerate(reader.pages):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            # 個別ページの抽出失敗はスキップ（次のページを試す）
            continue
        if not text:
            continue
        # 600字超は段落分割
        for j, piece in enumerate(_split_text(text)):
            suffix = f"p{i+1}" if j == 0 else f"p{i+1}-{j}"
            chunks.append(
                Chunk(chunk_id=f"{filename}#{suffix}", source=filename, text=piece)
            )
    return chunks


def _parse_xlsx(filename: str, content: bytes) -> list[Chunk]:
    """Excel: シート×行をチャンクに変換。1行=1チャンク（ヘッダー付き "key: value"）。"""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    chunks: list[Chunk] = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows_iter = ws.iter_rows(values_only=True)
        try:
            header = [str(c) if c is not None else f"col{i}" for i, c in enumerate(next(rows_iter))]
        except StopIteration:
            continue
        for r_idx, row in enumerate(rows_iter, start=2):
            parts = [
                f"{h}: {v}" for h, v in zip(header, row)
                if v is not None and str(v).strip()
            ]
            if not parts:
                continue
            text = " / ".join(parts)
            chunks.append(
                Chunk(chunk_id=f"{filename}#{sheet}!r{r_idx}", source=filename, text=text)
            )
    return chunks


def _parse_pptx(filename: str, content: bytes) -> list[Chunk]:
    """PowerPoint: スライド1枚を1チャンクに変換（タイトル + テキスト枠 + ノート）。"""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    chunks: list[Chunk] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
        # 発表者ノート
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(f"[ノート] {note}")
        if not parts:
            continue
        text = "\n".join(parts)
        chunks.append(
            Chunk(chunk_id=f"{filename}#slide{i}", source=filename, text=text)
        )
    return chunks


def parse(filename: str, content: bytes) -> list[Chunk]:
    filename = _safe_filename(filename)
    fmt = _detect_format(filename)
    if fmt in ("markdown", "text", "json"):
        text = content.decode("utf-8-sig", errors="replace")
        return _parse_text(filename, text)
    if fmt == "csv":
        return _parse_csv(filename, content)
    if fmt == "pdf":
        return _parse_pdf(filename, content)
    if fmt == "xlsx":
        return _parse_xlsx(filename, content)
    if fmt == "pptx":
        return _parse_pptx(filename, content)
    raise ValueError(f"unsupported format: {Path(filename).suffix}")


# =====================================================
# 検出
# =====================================================
def _scan_pii(text: str, industry: str = "general") -> dict[str, int]:
    """PII検出: 各ルールにマッチした件数を返す。"""
    counts: dict[str, int] = {}
    for rule in build_rules(industry):
        n = len(rule.pattern.findall(text))
        if n:
            counts[rule.name] = counts.get(rule.name, 0) + n
    return counts


def _scan_confidential(text: str) -> list[str]:
    return [m for m in CONFIDENTIAL_MARKERS if m in text]


def _scan_names(text: str) -> int:
    return len(NAME_HEURISTIC.findall(text))


def scan_findings(chunks: list[Chunk], industry: str = "general") -> FileFindings:
    findings = FileFindings()
    full_text = "\n\n".join(c.text for c in chunks)
    findings.pii_counts = _scan_pii(full_text, industry)
    findings.confidential_markers = _scan_confidential(full_text)
    findings.name_candidates = _scan_names(full_text)
    return findings


def scan_chunk_findings(chunks: list[Chunk], industry: str = "general") -> list[ChunkFindings]:
    """各チャンクごとに PII / 機密マーカー / 推奨判定を返す。"""
    out: list[ChunkFindings] = []
    for c in chunks:
        cf = ChunkFindings(
            pii_counts=_scan_pii(c.text, industry),
            confidential_markers=_scan_confidential(c.text),
            name_candidates=_scan_names(c.text),
        )
        cf.recommendation, cf.reason = assess(
            FileFindings(
                pii_counts=cf.pii_counts,
                confidential_markers=cf.confidential_markers,
                name_candidates=cf.name_candidates,
            ),
            n_chunks=1,
        )
        out.append(cf)
    return out


# =====================================================
# 推奨判定
# =====================================================
def assess(findings: FileFindings, n_chunks: int) -> tuple[Recommendation, str]:
    """マイナンバー・カード・氏名多数なら danger、PII軽微なら warn、それ以外 ok。"""
    pii = findings.pii_counts
    if pii.get("my_number", 0) >= 1:
        return "danger", f"マイナンバー疑い {pii['my_number']}件 — 法的リスクのため取り込み非推奨"
    if pii.get("credit_card", 0) >= 1:
        return "danger", f"クレジットカード番号 {pii['credit_card']}件 — PCI-DSS抵触可能性"
    if findings.name_candidates >= max(20, n_chunks * 2):
        return "danger", f"個人氏名候補 {findings.name_candidates}件 — 個人情報集の可能性"

    warnings: list[str] = []
    if findings.confidential_markers:
        warnings.append(f"機密マーカー検出: {', '.join(findings.confidential_markers)}")
    if findings.name_candidates >= 5:
        warnings.append(f"個人氏名候補 {findings.name_candidates}件 → マスキング推奨")
    pii_warn = {k: v for k, v in pii.items() if k in ("email", "phone_jp", "ip_address", "url")}
    if pii_warn:
        warnings.append("PII: " + ", ".join(f"{k}{v}件" for k, v in pii_warn.items()) + " → 自動マスク予定")

    if warnings:
        return "warn", " / ".join(warnings)
    return "ok", "懸念事項なし"


# =====================================================
# 全体パイプライン
# =====================================================
def analyze(filename: str, content: bytes, industry: str = "general") -> FileAnalysis:
    safe_name = _safe_filename(filename)
    chunks = parse(safe_name, content)
    findings = scan_findings(chunks, industry)
    chunk_findings = scan_chunk_findings(chunks, industry)
    rec, reason = assess(findings, len(chunks))
    # 内容が無い場合の特別ケース
    if not chunks:
        rec = "warn"
        reason = "テキストを抽出できませんでした（空ファイル / 画像のみPDF / 空 Excel 等の可能性）"
    return FileAnalysis(
        filename=safe_name,
        sha256=hashlib.sha256(content).hexdigest(),
        size_bytes=len(content),
        format=_detect_format(safe_name),
        chunks=chunks,
        findings=findings,
        recommendation=rec,
        reason=reason,
        chunk_findings=chunk_findings,
    )


def ingest(
    analysis: FileAnalysis,
    faq_master_dir: Path,
    apply_masking: bool = True,
    industry: str = "general",
    excluded_chunk_ids: set[str] | None = None,
) -> int:
    """マスク適用後にFAQマスターへ書き出し。

    Args:
      excluded_chunk_ids: スキップするチャンクIDの集合。指定された chunk_id は
                          書き出し対象から除外される（部分取り込み）
    """
    faq_master_dir.mkdir(parents=True, exist_ok=True)
    out_path = faq_master_dir / Path(analysis.filename).with_suffix(".md").name

    rules = build_rules(industry) if apply_masking else None
    excluded = excluded_chunk_ids or set()
    body_parts = [f"# {analysis.filename}\n\n_取り込み元: {analysis.sha256[:12]}_\n"]
    written = 0
    for chunk in analysis.chunks:
        if chunk.chunk_id in excluded:
            continue
        text = mask(chunk.text, rules) if rules is not None else chunk.text
        body_parts.append(f"\n## {chunk.chunk_id}\n\n{text}\n")
        written += 1
    if written == 0:
        # 全チャンク除外なら出力しない
        return 0
    out_path.write_text("\n".join(body_parts), encoding="utf-8")
    return written
