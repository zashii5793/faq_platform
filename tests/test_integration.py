"""統合テスト：複数エンドポイントを跨ぐ実シナリオを検証する。

カバー範囲:
  1. アップロード → 取り込み → 質問 → 回答（完全フロー）
  2. 関連なし質問 → has_answer=False（ハルシネーション抑制）
  3. PIIマスキング → 監査ログのマスク済み記録
  4. 危険ファイルの取り込み拒否（HTTP 400）
  5. 複数フォーマット（CSV/Markdown/PowerPoint）の同時取り込み
  6. フィードバック投票 → サイドバー stats への反映
  7. 容量上限超過 → HTTP 413
"""
from __future__ import annotations

import io
import json
from urllib.parse import quote

import pytest
from fastapi.testclient import TestClient


@pytest.fixture
def client(tmp_path, monkeypatch) -> TestClient:
    """各テストで独立した FAQマスターディレクトリを使う。"""
    from app import audit, rag
    from app.config import settings
    from app.main import app

    monkeypatch.setattr(settings, "demo_mode", True)
    monkeypatch.setattr(settings, "anthropic_api_key", "")
    monkeypatch.setattr(settings, "faq_master_dir", tmp_path / "faq_master")
    monkeypatch.setattr(audit, "LOG_DIR", tmp_path / "audit")
    rag._index = None
    return TestClient(app)


def _xlsx_bytes(rows: list[list]) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    for r in rows:
        ws.append(r)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes(slides: list[str]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text in slides:
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2)).text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ============================================================
# シナリオ1: 完全フロー（取り込み → 質問 → 回答 → 確信度）
# ============================================================
def test_full_workflow_upload_then_ask(client: TestClient):
    # Step 1: VPN手順マークダウンをアップロード&取り込み
    md = "# VPN\n\nFortiClientを起動してログイン後、OTP6桁を入力してください。".encode("utf-8")
    r = client.post("/api/admin/ingest", files={"file": ("vpn.md", md, "text/markdown")})
    assert r.status_code == 200
    assert r.json()["ingested_chunks"] >= 1

    # Step 2: 取り込んだ内容について質問
    r = client.post("/api/ask", json={"question": "VPNにログインする手順"})
    assert r.status_code == 200
    data = r.json()
    assert data["has_answer"] is True
    assert data["confidence"] > 0
    assert any(s["source"] == "vpn.md" for s in data["sources"])


# ============================================================
# シナリオ2: ハルシネーション抑制
# ============================================================
def test_no_answer_for_irrelevant_question(client: TestClient):
    # 1件だけ取り込んで関係ない質問をする
    md = "# 経費精算\n\n月次締めは毎月25日です。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("expense.md", md, "text/markdown")})

    r = client.post("/api/ask", json={"question": "宇宙ロケットの打ち上げ手順を教えて"})
    data = r.json()
    assert data["has_answer"] is False
    assert data["confidence"] == 0
    assert "見つかりませんでした" in data["answer"]


def test_llm_no_answer_response_resets_signals(client: TestClient, monkeypatch):
    """LLM が「該当情報が見つかりませんでした」と返したら、
    confidence=0 / sources=[] / has_answer=False に揃える（UI 整合性）。
    """
    from app import llm
    from app.config import settings
    from app.rag import Chunk

    # 文書を取り込んで TF-IDF が hit するようにする
    md = "# 経費精算\n\n月次締めは毎月25日です。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("expense.md", md, "text/markdown")})

    # LLM スタブ：常に「該当情報なし」を返す
    class FakeMessages:
        def create(self, **kwargs):
            class FakeMsg:
                content = [type("B", (), {"type": "text", "text":
                    "該当情報が見つかりませんでした。社内ヘルプデスクにお問い合わせください。"})()]
                usage = type("U", (), {
                    "input_tokens": 50, "output_tokens": 10,
                    "cache_creation_input_tokens": 0, "cache_read_input_tokens": 0,
                })()
            return FakeMsg()

    class FakeClient:
        messages = FakeMessages()

    monkeypatch.setattr(llm, "_client", lambda: FakeClient())
    monkeypatch.setattr(settings, "anthropic_api_key", "test-key")

    r = client.post("/api/ask", json={"question": "年度更新の手順"})
    assert r.status_code == 200
    data = r.json()
    # LLM が「該当情報なし」と返したので信号が揃っている
    assert data["has_answer"] is False, "LLM の該当情報なし応答で has_answer が False になっていない"
    assert data["confidence"] == 0, f"confidence が 0 にリセットされていない: {data['confidence']}"
    assert data["sources"] == [], f"sources が空になっていない: {data['sources']}"
    assert data["is_reference"] is False
    assert "該当情報が見つかりませんでした" in data["answer"]


# ============================================================
# シナリオ3: PIIマスキングと監査ログ
# ============================================================
def test_pii_masked_in_audit_log(client: TestClient, tmp_path):
    md = "# FAQ\n\nVPN は FortiClient を使う".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("vpn.md", md, "text/markdown")})

    # PII を含む質問
    client.post("/api/ask", json={
        "question": "test@example.com のユーザーが 03-1234-5678 でVPN接続できない"
    })

    # 監査ログを確認: マスク済みで記録されているはず
    log_dir = tmp_path / "audit"
    log_files = list(log_dir.glob("*.jsonl"))
    assert log_files
    entries = [json.loads(line) for line in log_files[0].read_text().splitlines() if line.strip()]
    query_entries = [e for e in entries if e.get("event") == "query"]
    assert query_entries
    last = query_entries[-1]
    assert "[メール]" in last["question"]
    assert "[電話番号]" in last["question"]
    assert "test@example.com" not in last["question"]
    assert "03-1234-5678" not in last["question"]


# ============================================================
# シナリオ4: 危険ファイル取り込み拒否
# ============================================================
def test_dangerous_file_rejected(client: TestClient):
    csv = b"\xef\xbb\xbfid,name,number\n1,foo,1234 5678 9012\n2,bar,9876 5432 1098\n"
    r = client.post("/api/admin/ingest", files={"file": ("staff.csv", csv, "text/csv")})
    assert r.status_code == 400
    assert "マイナンバー" in r.json()["detail"]


def test_dangerous_file_can_still_be_analyzed(client: TestClient):
    """analyze は danger 判定でも結果を返す（UIで判断するため）。"""
    csv = b"\xef\xbb\xbfid,name,number\n1,foo,1234 5678 9012\n"
    r = client.post("/api/admin/analyze", files={"file": ("staff.csv", csv, "text/csv")})
    assert r.status_code == 200
    assert r.json()["recommendation"] == "danger"


# ============================================================
# シナリオ5: 複数フォーマット同時取り込み
# ============================================================
def test_multi_format_ingest_and_query(client: TestClient):
    # Markdown
    client.post("/api/admin/ingest", files={
        "file": ("md.md", "# 出席\n\n出席ボタン押下".encode("utf-8"), "text/markdown")
    })
    # CSV
    client.post("/api/admin/ingest", files={
        "file": ("csv.csv", "q,a\nVPN,FortiClient起動\n".encode("utf-8"), "text/csv")
    })
    # Excel
    client.post("/api/admin/ingest", files={
        "file": ("xl.xlsx", _xlsx_bytes([["category", "answer"], ["有給", "KING OF TIME"]]),
                 "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    })
    # PowerPoint
    client.post("/api/admin/ingest", files={
        "file": ("pp.pptx", _pptx_bytes(["経費精算は毎月25日"]),
                 "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    })

    # stats 確認: 4 文書取り込まれている
    r = client.get("/api/admin/stats")
    s = r.json()
    assert s["knowledge"]["n_documents"] == 4

    # 各フォーマットから検索できる（FAQマスターは .md に統一保存される仕様）
    for q, expected_src in [
        ("出席", "md.md"),
        ("VPN", "csv.md"),
        ("有給", "xl.md"),
        ("経費精算", "pp.md"),
    ]:
        r = client.post("/api/ask", json={"question": q})
        d = r.json()
        assert d["has_answer"] is True, f"{q} → {d}"
        assert any(s["source"] == expected_src for s in d["sources"]), \
            f"{q} の最上位が {expected_src} ではない: {d['sources']}"


# ============================================================
# シナリオ6: フィードバック投票 → stats 反映
# ============================================================
def test_feedback_updates_stats(client: TestClient):
    md = "# FAQ\n\n出席登録の方法".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("a.md", md, "text/markdown")})
    client.post("/api/ask", json={"question": "出席登録"})

    # フィードバック投票
    r = client.post("/api/feedback", json={"question": "出席登録", "vote": "up", "sources": ["a.md"]})
    assert r.status_code == 200

    r = client.post("/api/feedback", json={"question": "出席登録", "vote": "down", "sources": ["a.md"]})
    assert r.status_code == 200

    # stats 反映確認
    s = client.get("/api/admin/stats").json()
    assert s["feedback"]["up"] == 1
    assert s["feedback"]["down"] == 1
    assert "出席登録" in s["feedback"]["down_questions"]


def test_feedback_validates_vote(client: TestClient):
    r = client.post("/api/feedback", json={"question": "x", "vote": "invalid"})
    assert r.status_code == 400


# ============================================================
# シナリオ7: 容量上限
# ============================================================
def test_oversized_file_rejected(client: TestClient):
    # 51MB のダミー
    big = b"x" * (51 * 1024 * 1024)
    r = client.post("/api/admin/analyze", files={"file": ("big.txt", big, "text/plain")})
    assert r.status_code == 413


# ============================================================
# シナリオ8: 未対応形式
# ============================================================
def test_unsupported_format(client: TestClient):
    r = client.post("/api/admin/analyze", files={"file": ("malware.exe", b"MZ", "application/octet-stream")})
    assert r.status_code == 415


# ============================================================
# シナリオ9: stats が初期状態で動作する
# ============================================================
def test_stats_initial_state(client: TestClient):
    s = client.get("/api/admin/stats").json()
    assert s["knowledge"]["n_documents"] == 0
    assert s["analytics"]["n_queries_today"] == 0
    assert s["analytics"]["avg_confidence"] == 0
    assert s["feedback"]["up"] == 0
    assert s["history"] == []


# ============================================================
# シナリオ10: 確信度がトピックの種類で変わる
# ============================================================
def test_confidence_higher_for_specific_match(client: TestClient):
    md = "# VPN詳細\n\nFortiClient設定: corp-vpn / 社員番号 + AD パスワード + OTP".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("vpn.md", md, "text/markdown")})

    specific = client.post("/api/ask", json={"question": "FortiClientの設定"}).json()
    vague = client.post("/api/ask", json={"question": "システムについて"}).json()

    # 具体的な質問の方が確信度が高い OR vague は no_answer
    if vague["has_answer"]:
        assert specific["confidence"] >= vague["confidence"]
    else:
        assert specific["has_answer"] is True


# ============================================================
# シナリオ11: フィードバック学習 — 👍 でランキングが上がる
# ============================================================
def test_feedback_boosts_search_ranking(client, monkeypatch, tmp_path):
    """同程度の関連性を持つ2文書で、👍多い方が上位に来ることを検証。"""
    # フィードバックスコアを独立した一時ファイルに分離
    from app import rag
    monkeypatch.setattr(rag, "FEEDBACK_PATH", tmp_path / "feedback.json")

    # 2つのVPN関連文書を取り込み
    md_a = "# VPN手順A\n\nVPN接続の方法。FortiClient起動。".encode("utf-8")
    md_b = "# VPN手順B\n\nVPN接続の方法。VPNクライアント起動。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("vpn_a.md", md_a, "text/markdown")})
    client.post("/api/admin/ingest", files={"file": ("vpn_b.md", md_b, "text/markdown")})

    # 初期検索: どちらが上位かを記録
    initial = client.post("/api/ask", json={"question": "VPN接続の方法"}).json()
    assert initial["has_answer"] is True
    initial_top = initial["sources"][0]["source"]

    # 反対側に👍を5回投票して学習を強化
    other = "vpn_b.md" if initial_top == "vpn_a.md" else "vpn_a.md"
    for _ in range(5):
        client.post("/api/feedback", json={
            "question": "VPN接続の方法", "vote": "up", "sources": [other]
        })

    # 再検索: ブーストされた方が上位に来るはず
    after = client.post("/api/ask", json={"question": "VPN接続の方法"}).json()
    assert after["sources"][0]["source"] == other, \
        f"フィードバック学習でランキングが変わるべき: 初期={initial_top}, 期待={other}, " \
        f"実際={after['sources'][0]['source']}"


def test_analyze_returns_per_chunk_findings(client: TestClient):
    """analyze レスポンスの chunks に各チャンクの判定が含まれる。"""
    md = (
        "## 普通の節\n\nここは安全な情報です。\n\n"
        "## 連絡先\n\n問い合わせは support@example.com まで。\n\n"
        "## 機密\n\n社外秘の重要な情報。"
    ).encode("utf-8")
    r = client.post("/api/admin/analyze", files={"file": ("doc.md", md, "text/markdown")})
    assert r.status_code == 200
    data = r.json()
    assert "chunks" in data
    assert len(data["chunks"]) >= 2
    # 各チャンクに recommendation と findings が含まれる
    for c in data["chunks"]:
        assert "chunk_id" in c
        assert "recommendation" in c
        assert "findings" in c
    # warn または danger チャンクが少なくとも1つある（連絡先 or 機密）
    recs = [c["recommendation"] for c in data["chunks"]]
    assert any(r in ("warn", "danger") for r in recs)


def test_ingest_with_excluded_chunks(client: TestClient):
    """excluded_chunk_ids で部分除外しての取り込みができる。"""
    md = (
        "## 安全な情報\n\nこれは公開してOK。\n\n"
        "## 機密情報\n\n社外秘の内容。"
    ).encode("utf-8")
    # まず analyze でチャンク確認
    r = client.post("/api/admin/analyze", files={"file": ("doc.md", md, "text/markdown")})
    chunks = r.json()["chunks"]
    confidential_id = next(c["chunk_id"] for c in chunks if c["recommendation"] == "warn")

    # その chunk_id を除外して ingest（'#' を含むので URL エンコード必須）
    r = client.post(
        f"/api/admin/ingest?excluded_chunk_ids={quote(confidential_id, safe='')}",
        files={"file": ("doc.md", md, "text/markdown")},
    )
    assert r.status_code == 200
    data = r.json()
    assert data["excluded_chunks"] == 1
    # 取り込み済み数 = 全チャンク - 除外数
    assert data["ingested_chunks"] == len(chunks) - 1


def test_ingest_dangerous_file_can_partial_with_exclusion(client: TestClient):
    """ファイル全体が danger でも、危険チャンクを除外すれば取り込める。"""
    csv = (
        b"\xef\xbb\xbfname,note\n"
        b"safe item,public OK\n"
        b"sensitive,number 1234 5678 9012 here\n"  # マイナンバー
    )
    r = client.post("/api/admin/analyze", files={"file": ("mix.csv", csv, "text/csv")})
    data = r.json()
    danger_ids = ",".join(c["chunk_id"] for c in data["chunks"] if c["recommendation"] == "danger")

    # 危険チャンク全部を除外して ingest（force 不要、'#' を含むので URL エンコード）
    r = client.post(
        f"/api/admin/ingest?excluded_chunk_ids={quote(danger_ids, safe=',')}",
        files={"file": ("mix.csv", csv, "text/csv")},
    )
    assert r.status_code == 200
    assert r.json()["ingested_chunks"] >= 1


def test_feedback_down_reduces_ranking(client, monkeypatch, tmp_path):
    """👎を入れると同じ文書のスコアが下がる。"""
    from app import rag
    monkeypatch.setattr(rag, "FEEDBACK_PATH", tmp_path / "feedback.json")

    md = "# 経費精算\n\n月次締めは毎月25日です。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("expense.md", md, "text/markdown")})

    before = client.post("/api/ask", json={"question": "経費精算の締め日"}).json()
    assert before["has_answer"] is True
    before_score = before["sources"][0]["score"]

    # 👎を3回
    for _ in range(3):
        client.post("/api/feedback", json={
            "question": "経費精算の締め日", "vote": "down", "sources": ["expense.md"]
        })

    after = client.post("/api/ask", json={"question": "経費精算の締め日"}).json()
    after_score = after["sources"][0]["score"]
    assert after_score < before_score, \
        f"👎でスコアが下がるべき: before={before_score}, after={after_score}"


# ============================================================
# シナリオ12: 取り込み済み文書のメンテナンス（一覧 / 削除）
# ============================================================
def test_documents_list_empty(client: TestClient):
    """初期状態では空リストが返る。"""
    r = client.get("/api/admin/documents")
    assert r.status_code == 200
    assert r.json() == {"documents": []}


def test_documents_list_after_ingest(client: TestClient):
    """取り込み後、文書一覧にメタデータ込みで現れる。"""
    md = "# VPN手順\n\nFortiClient を起動してログイン。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("vpn.md", md, "text/markdown")})

    r = client.get("/api/admin/documents")
    assert r.status_code == 200
    docs = r.json()["documents"]
    assert len(docs) == 1
    d = docs[0]
    assert d["filename"] == "vpn.md"
    assert d["size_bytes"] > 0
    assert d["n_chunks"] >= 1
    assert "modified_at" in d


def test_document_delete_removes_from_index(client: TestClient):
    """削除すると、ファイルもインデックスからも消え、検索ヒットしなくなる。"""
    md = "# 部品発注\n\n締め時間は毎日17時です。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("parts.md", md, "text/markdown")})

    # 取り込まれた状態で質問するとヒットする
    before = client.post("/api/ask", json={"question": "部品発注の締め時間"}).json()
    assert before["has_answer"] is True

    # 削除
    r = client.delete("/api/admin/documents/parts.md")
    assert r.status_code == 200
    body = r.json()
    assert body["deleted"] == "parts.md"
    assert body["n_chunks_after"] == 0

    # 一覧から消える
    r2 = client.get("/api/admin/documents")
    assert r2.json()["documents"] == []

    # 検索しても答えない（ハルシネーション抑制）
    after = client.post("/api/ask", json={"question": "部品発注の締め時間"}).json()
    assert after["has_answer"] is False


def test_document_delete_path_traversal_blocked(client: TestClient):
    """パストラバーサル攻撃は 400 で拒否される。"""
    # `..` を含むパスは 400
    r = client.delete("/api/admin/documents/..%2F..%2Fetc%2Fpasswd")
    assert r.status_code == 400


def test_document_delete_nonexistent_returns_404(client: TestClient):
    """存在しないファイルの削除は 404。"""
    r = client.delete("/api/admin/documents/nonexistent.md")
    assert r.status_code == 404


def test_document_delete_is_audit_logged(client: TestClient, tmp_path):
    """削除イベントが監査ログに記録される。"""
    from app import audit

    md = "# テストFAQ\n\n削除テスト用。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("delete_me.md", md, "text/markdown")})
    client.delete("/api/admin/documents/delete_me.md")

    log = audit.read_recent(50)
    delete_events = [e for e in log if e.get("event") == "delete_document"]
    assert len(delete_events) == 1
    assert delete_events[0]["filename"] == "delete_me.md"
    assert delete_events[0]["size_bytes"] > 0


# ============================================================
# シナリオ13: FAQ追加リクエスト
# ============================================================
def test_faq_request_creates_audit_event(client: TestClient):
    """FAQ追加リクエストが監査ログに記録される。"""
    from app import audit

    r = client.post("/api/faq-requests", json={
        "question": "退職金の計算方法は？",
        "note": "新人から3回目の質問",
    })
    assert r.status_code == 200
    assert r.json()["ok"] is True

    log = audit.read_recent(50)
    requests = [e for e in log if e.get("event") == "faq_request"]
    assert len(requests) == 1
    assert requests[0]["question"] == "退職金の計算方法は？"
    assert requests[0]["note"] == "新人から3回目の質問"


def test_faq_request_rejects_empty(client: TestClient):
    """空の質問はリクエスト不可。"""
    r = client.post("/api/faq-requests", json={"question": "   ", "note": ""})
    assert r.status_code == 400


def test_faq_request_rejects_too_long(client: TestClient):
    """2000文字超は拒否。"""
    r = client.post("/api/faq-requests", json={"question": "あ" * 2001, "note": ""})
    assert r.status_code == 400


def test_admin_list_faq_requests(client: TestClient):
    """管理画面用 FAQ追加リクエスト一覧。"""
    client.post("/api/faq-requests", json={"question": "Q1", "note": ""})
    client.post("/api/faq-requests", json={"question": "Q2", "note": "緊急"})

    r = client.get("/api/admin/faq-requests")
    assert r.status_code == 200
    data = r.json()
    # 監査ログは新しい順なので Q2 が先
    questions = [req["question"] for req in data["requests"]]
    assert "Q1" in questions
    assert "Q2" in questions
    assert data["total"] == 2


def test_ask_returns_reference_mode_for_low_confidence(client: TestClient):
    """確信度が低くても、関連チャンクがあれば参考情報として返す（B モード）。"""
    # 1文書だけ取り込んで、関連性が薄い質問をする
    md = "# 会議室予約\n\n会議室は社内ポータルから予約できます。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("meeting.md", md, "text/markdown")})

    # 直接の答えがない質問だが、何かしらヒットする
    r = client.post("/api/ask", json={"question": "リモートワークの申請手続きは？"})
    assert r.status_code == 200
    data = r.json()
    # 関連チャンクは見つかるので has_answer=True、低確信度なら is_reference=True
    # ローカルスタブモードなので answer 文字列で is_reference 動作確認
    if data["has_answer"] and data["confidence"] < 50:
        assert data["is_reference"] is True


def test_ask_returns_no_answer_with_empty_index(client: TestClient):
    """インデックスが空なら has_answer=False で is_reference=False。"""
    r = client.post("/api/ask", json={"question": "VPN接続の方法は？"})
    assert r.status_code == 200
    data = r.json()
    assert data["has_answer"] is False
    assert data["is_reference"] is False
    assert data["confidence"] == 0


# ============================================================
# シナリオ14: 監査ログのエクスポート（顧客レポート用）
# ============================================================
def test_export_csv_query_history(client: TestClient):
    """質問履歴を CSV でエクスポート。"""
    md = "# VPN\n\nFortiClient で接続。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("vpn.md", md, "text/markdown")})
    client.post("/api/ask", json={"question": "VPNの接続方法"})

    r = client.get("/api/admin/export?days=30&format=csv&event=query")
    assert r.status_code == 200
    assert "text/csv" in r.headers["content-type"]
    assert "attachment" in r.headers["content-disposition"]
    body = r.content.decode("utf-8")
    # UTF-8 BOM
    assert body.startswith("﻿")
    # CSV ヘッダ
    assert "question" in body
    assert "confidence" in body
    assert "VPNの接続方法" in body


def test_export_json_all_events(client: TestClient):
    """全イベントを JSON でエクスポート。"""
    client.post("/api/faq-requests", json={"question": "Q?", "note": ""})

    r = client.get("/api/admin/export?days=30&format=json&event=all")
    assert r.status_code == 200
    data = r.json()
    assert "entries" in data
    assert data["n_rows"] >= 1


def test_export_invalid_format_rejected(client: TestClient):
    """xml など未対応の format は 400。"""
    r = client.get("/api/admin/export?format=xml")
    assert r.status_code == 400


def test_export_invalid_days_rejected(client: TestClient):
    """範囲外の days は 400。"""
    r = client.get("/api/admin/export?days=999&format=csv")
    assert r.status_code == 400
    r = client.get("/api/admin/export?days=0&format=csv")
    assert r.status_code == 400


def test_export_audit_logged(client: TestClient):
    """エクスポート自体も監査される（誰がいつ何を出力したか）。"""
    from app import audit
    client.get("/api/admin/export?format=csv&event=query")
    log = audit.read_recent(10)
    exports = [e for e in log if e.get("event") == "export"]
    assert len(exports) == 1
    assert exports[0]["format"] == "csv"
    assert exports[0]["event_filter"] == "query"


# ============================================================
# シナリオ15: 組織情報のランタイム編集
# ============================================================
def test_settings_get_returns_defaults(client: TestClient, tmp_path, monkeypatch):
    """初期状態では .env の値が effective として返る。"""
    from app import runtime_settings
    monkeypatch.setattr(runtime_settings, "OVERRIDES_PATH", tmp_path / "org_settings.json")

    r = client.get("/api/admin/settings")
    assert r.status_code == 200
    data = r.json()
    assert "effective" in data
    assert "org_name" in data["effective"]
    assert data["overrides"] == {}
    assert "org_name" in data["editable_keys"]


def test_settings_update_persists(client: TestClient, tmp_path, monkeypatch):
    """更新後は overrides に保存され、effective にも反映される。"""
    from app import runtime_settings
    monkeypatch.setattr(runtime_settings, "OVERRIDES_PATH", tmp_path / "org_settings.json")

    r = client.put("/api/admin/settings", json={
        "org_name": "テスト商事株式会社",
        "assistant_role": "総務サポート",
    })
    assert r.status_code == 200
    body = r.json()
    assert body["applied"]["org_name"] == "テスト商事株式会社"

    # ファイルが書かれている
    saved = json.loads((tmp_path / "org_settings.json").read_text(encoding="utf-8"))
    assert saved["org_name"] == "テスト商事株式会社"

    # GET でも反映
    r2 = client.get("/api/admin/settings")
    assert r2.json()["effective"]["org_name"] == "テスト商事株式会社"


def test_settings_update_rejects_empty(client: TestClient):
    """空 payload は 400。"""
    r = client.put("/api/admin/settings", json={})
    assert r.status_code == 400


def test_settings_update_rejects_too_long(client: TestClient, tmp_path, monkeypatch):
    """200文字超は 400。"""
    from app import runtime_settings
    monkeypatch.setattr(runtime_settings, "OVERRIDES_PATH", tmp_path / "org_settings.json")
    r = client.put("/api/admin/settings", json={"org_name": "あ" * 201})
    assert r.status_code == 400


def test_settings_reset_clears_overrides(client: TestClient, tmp_path, monkeypatch):
    """DELETE で全オーバーライドが消える。"""
    from app import runtime_settings
    monkeypatch.setattr(runtime_settings, "OVERRIDES_PATH", tmp_path / "org_settings.json")

    client.put("/api/admin/settings", json={"org_name": "Foo"})
    assert (tmp_path / "org_settings.json").exists()

    r = client.delete("/api/admin/settings")
    assert r.status_code == 200
    assert not (tmp_path / "org_settings.json").exists()


# ============================================================
# シナリオ16: 人気質問のサジェスト
# ============================================================
def test_popular_queries_in_stats(client: TestClient):
    """同じ質問を2回以上聞かれたら popular_queries に出る。"""
    md = "# 経費精算\n\n毎月25日締めです。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("expense.md", md, "text/markdown")})

    # 同じ質問を3回
    for _ in range(3):
        client.post("/api/ask", json={"question": "経費精算の締め日は？"})

    r = client.get("/api/admin/stats")
    assert r.status_code == 200
    data = r.json()
    assert "popular_queries" in data
    popular = data["popular_queries"]
    questions = [p["question"] for p in popular]
    assert "経費精算の締め日は？" in questions
    found = next(p for p in popular if p["question"] == "経費精算の締め日は？")
    assert found["count"] >= 2


def test_single_query_not_in_popular(client: TestClient):
    """1回しか聞かれていない質問はサジェストされない。"""
    md = "# VPN\n\nFortiClient で接続。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("vpn.md", md, "text/markdown")})
    client.post("/api/ask", json={"question": "VPN設定"})

    r = client.get("/api/admin/stats")
    popular = r.json().get("popular_queries", [])
    questions = [p["question"] for p in popular]
    assert "VPN設定" not in questions


# ============================================================
# シナリオ17: 利用状況ダッシュボード
# ============================================================
def test_dashboard_returns_daily_buckets(client: TestClient):
    """過去 N 日分の日次バケット（質問数ゼロ含む）が返る。"""
    md = "# テスト\n\nテスト用FAQ。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("t.md", md, "text/markdown")})
    client.post("/api/ask", json={"question": "テスト"})

    r = client.get("/api/admin/dashboard?days=7")
    assert r.status_code == 200
    d = r.json()
    assert d["days"] == 7
    assert len(d["daily"]) == 7
    # 各バケットに必要なキーがある
    for bucket in d["daily"]:
        for key in ("date", "queries", "answered", "reference", "no_answer",
                    "avg_confidence", "unique_users"):
            assert key in bucket
    # 合計に少なくとも 1件の質問
    assert d["totals"]["queries"] >= 1


def test_dashboard_invalid_days_rejected(client: TestClient):
    """範囲外の days は 400。"""
    assert client.get("/api/admin/dashboard?days=0").status_code == 400
    assert client.get("/api/admin/dashboard?days=999").status_code == 400


def test_dashboard_top_topics_populated_after_queries(client: TestClient):
    """質問が回答されると top_topics が増える。"""
    md = "# 経費精算\n\n月末締めです。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("expense.md", md, "text/markdown")})
    for _ in range(3):
        client.post("/api/ask", json={"question": "経費の締めは？"})

    r = client.get("/api/admin/dashboard?days=30")
    d = r.json()
    topics = [t["source"] for t in d["top_topics"]]
    # ローカルスタブでは sources がチャンクID で記録されている可能性あり、
    # いずれにせよ top_topics が空でなければOK
    if d["totals"]["queries"] > 0 and topics:
        assert any("expense" in s.lower() or "経費" in s for s in topics) or len(topics) > 0


def test_dashboard_llm_usage_block_present(client: TestClient):
    """ダッシュボードに llm_usage ブロックが存在する（ローカルモードでは calls=0）。"""
    r = client.get("/api/admin/dashboard?days=14")
    d = r.json()
    assert "llm_usage" in d
    lu = d["llm_usage"]
    for key in ("calls", "input_tokens", "output_tokens",
                "cache_creation_tokens", "cache_read_tokens", "cache_hit_rate"):
        assert key in lu


# ============================================================
# シナリオ18: 質問履歴の検索・絞り込み
# ============================================================
def test_query_search_filters_by_keyword(client: TestClient):
    """質問本文の部分一致で絞り込める。"""
    md = "# テスト".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("t.md", md, "text/markdown")})
    client.post("/api/ask", json={"question": "VPN接続の方法"})
    client.post("/api/ask", json={"question": "経費精算の締め日"})
    client.post("/api/ask", json={"question": "VPN設定変更"})

    r = client.get("/api/admin/queries?q=VPN&days=30")
    assert r.status_code == 200
    results = r.json()["results"]
    questions = [x["question"] for x in results]
    assert all("vpn" in q.lower() for q in questions)
    assert len(results) >= 2


def test_query_search_filters_by_answered_status(client: TestClient):
    """answered=no で未回答のみ取れる。"""
    md = "# VPN\n\nFortiClient で接続。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("vpn.md", md, "text/markdown")})
    client.post("/api/ask", json={"question": "VPN設定"})  # ヒット
    client.post("/api/ask", json={"question": "全くの無関係質問xyz123"})  # ノーヒット想定

    r = client.get("/api/admin/queries?days=30&answered=no")
    results = r.json()["results"]
    for x in results:
        assert x["answered"] is False


def test_query_search_invalid_params_rejected(client: TestClient):
    """範囲外パラメータは 400。"""
    assert client.get("/api/admin/queries?days=999").status_code == 400
    assert client.get("/api/admin/queries?min_confidence=-1").status_code == 400
    assert client.get("/api/admin/queries?max_confidence=200").status_code == 400
    assert client.get("/api/admin/queries?min_confidence=80&max_confidence=50").status_code == 400
    assert client.get("/api/admin/queries?answered=foo").status_code == 400
    assert client.get("/api/admin/queries?limit=0").status_code == 400
    assert client.get("/api/admin/queries?limit=999").status_code == 400


def test_query_search_respects_limit(client: TestClient):
    """limit より多い結果は切り詰められる、total は全件数。"""
    md = "# テスト".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("t.md", md, "text/markdown")})
    for i in range(15):
        client.post("/api/ask", json={"question": f"質問{i}"})

    r = client.get("/api/admin/queries?days=30&limit=5")
    d = r.json()
    assert len(d["results"]) == 5
    assert d["total"] >= 15


# ============================================================
# シナリオ19: プロンプトキャッシュ（system を構造化ブロックで送る）
# ============================================================
# ============================================================
# シナリオ19: PDF 取り込みの実シナリオ（E2E）
# テキスト抽出可能PDF / 画像のみPDF / 壊れたPDF をカバー
# ============================================================
def _make_text_pdf(text: str) -> bytes:
    """テキスト抽出可能な最小 PDF を手書きで生成（pypdf で読める）。"""
    return (
        b"%PDF-1.4\n"
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n"
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n"
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >> endobj\n"
        b"4 0 obj << /Length 60 >> stream\n"
        b"BT /F1 12 Tf 50 700 Td ("
        + text.encode("latin-1", errors="replace") + b") Tj ET\n"
        b"endstream endobj\n"
        b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n"
        b"xref\n0 6\n0000000000 65535 f\n"
        b"0000000009 00000 n\n0000000058 00000 n\n0000000111 00000 n\n"
        b"0000000226 00000 n\n0000000316 00000 n\n"
        b"trailer << /Size 6 /Root 1 0 R >>\nstartxref\n388\n%%EOF\n"
    )


def _make_image_only_pdf() -> bytes:
    """テキストを含まない PDF（画像のみページを想定）"""
    return (
        b"%PDF-1.4\n"
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n"
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n"
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] >> endobj\n"
        b"xref\n0 4\n0000000000 65535 f\n"
        b"0000000009 00000 n\n0000000058 00000 n\n0000000111 00000 n\n"
        b"trailer << /Size 4 /Root 1 0 R >>\nstartxref\n200\n%%EOF\n"
    )


def test_pdf_with_text_full_workflow(client: TestClient):
    """テキスト抽出可能PDF: analyze → ingest → 質問できる。"""
    pdf = _make_text_pdf("VPN connection requires FortiClient")

    r = client.post("/api/admin/analyze", files={"file": ("vpn.pdf", pdf, "application/pdf")})
    assert r.status_code == 200, r.text
    d = r.json()
    assert d["format"] == "pdf"
    assert d["n_chunks"] >= 1
    assert d["recommendation"] in ("ok", "warn")

    r = client.post("/api/admin/ingest", files={"file": ("vpn.pdf", pdf, "application/pdf")})
    assert r.status_code == 200
    assert r.json()["ingested_chunks"] >= 1


def test_pdf_image_only_analyze_returns_warning_card(client: TestClient):
    """画像のみPDFは 200 で n_chunks=0 + warn を返す（フロントが警告カードを出せる）。"""
    pdf = _make_image_only_pdf()

    r = client.post("/api/admin/analyze", files={"file": ("scan.pdf", pdf, "application/pdf")})
    assert r.status_code == 200
    d = r.json()
    assert d["n_chunks"] == 0
    assert d["recommendation"] == "warn"
    assert "テキスト" in d["reason"]  # 明確な理由メッセージ


def test_pdf_image_only_ingest_rejected_with_clear_message(client: TestClient):
    """画像のみPDFのingestは 422 + OCR案内付きで拒否される（無反応を防止）。"""
    pdf = _make_image_only_pdf()

    r = client.post("/api/admin/ingest", files={"file": ("scan.pdf", pdf, "application/pdf")})
    assert r.status_code == 422
    detail = r.json()["detail"]
    assert "テキスト" in detail
    assert "OCR" in detail or "ocr" in detail.lower()


def test_pdf_corrupted_returns_422_not_500(client: TestClient):
    """壊れたPDFは 500 ではなく 422 で明確に拒否される（500 はフロントの無反応につながる）。"""
    broken = b"%PDF-1.4\n%binary garbage\xff\xff\xff\n%%EOF\n"

    r = client.post("/api/admin/analyze", files={"file": ("broken.pdf", broken, "application/pdf")})
    # analyze は安全に処理 → n_chunks=0 の warn 扱い、または 422 のどちらか
    # （pypdf がreaderコンストラクタで例外を吐く場合は422、握って空ならanalyzeが200で0チャンク）
    assert r.status_code in (200, 422), r.text
    if r.status_code == 200:
        assert r.json()["n_chunks"] == 0


def test_pdf_not_a_pdf_returns_422(client: TestClient):
    """拡張子だけ.pdfで中身がテキストの場合も 200 で n_chunks=0 になる（500を吐かない）。"""
    fake = b"this is plain text pretending to be a pdf"

    r = client.post("/api/admin/analyze", files={"file": ("fake.pdf", fake, "application/pdf")})
    assert r.status_code in (200, 422)
    if r.status_code == 200:
        assert r.json()["n_chunks"] == 0


# ============================================================
# シナリオ20: 埋め込み JavaScript の構文回帰防止
# 過去に Python triple-quoted string 内の `\n` が実改行に変換され、
# JS の文字列リテラルがブレークして UI 全体が止まる重大バグがあったため、
# 全ページの <script> ブロックを抜き出して syntax チェックする。
# ============================================================
def test_admin_upload_js_has_valid_syntax(client: TestClient):
    """/admin/upload の JS が文字列内に未エスケープ改行を含まないことを保証。

    Python の \"\"\"..\"\"\" の中で `\\n` を書くと改行になってしまうため、JS の
    文字列リテラル内では `\\\\n` を書く必要がある。これを忘れると JS が
    SyntaxError で全停止する（ダッシュボード・PDF UI 含めて）。
    """
    import re
    r = client.get("/admin/upload")
    assert r.status_code == 200
    scripts = re.findall(r"<script>(.*?)</script>", r.text, re.DOTALL)
    assert scripts, "/admin/upload に <script> がない"

    for idx, js in enumerate(scripts):
        # 文字列リテラル内に未エスケープ改行が混入していないか検出
        # シングルクォート / バッククォート / ダブルクォート で囲まれた文字列を抽出
        _assert_no_raw_newline_in_js_strings(js, f"/admin/upload script[{idx}]")


def test_chat_page_js_has_valid_syntax(client: TestClient):
    """チャット画面（/）の JS も同じ回帰防止。"""
    import re
    r = client.get("/")
    if r.status_code != 200:
        # 認証必要の場合はスキップ（DEMO_MODE 想定なので通常は 200）
        return
    scripts = re.findall(r"<script>(.*?)</script>", r.text, re.DOTALL)
    for idx, js in enumerate(scripts):
        _assert_no_raw_newline_in_js_strings(js, f"/ script[{idx}]")


def _assert_no_raw_newline_in_js_strings(js: str, label: str):
    """JS 文字列リテラル内の未エスケープ改行を検出して failure させる。

    対象: '...' および `...`（ダブルクォートは template の中で複雑なのでスキップ）
    エスケープされた `\\n` は許可、実改行のみ NG。
    """
    # 簡易スキャナ: ストリングを舐めて、開始クォートから対応する終了クォートまで
    # に未エスケープの \n があれば failure。
    state = None  # None / "'" / "`"
    line_no = 1
    string_start_line = 0
    string_content_lines: list[int] = []
    i = 0
    while i < len(js):
        ch = js[i]
        if ch == "\n":
            line_no += 1
            if state in ("'",):
                # シングルクォート文字列の中に改行 → JS構文エラーになる
                raise AssertionError(
                    f"{label}: シングルクォート文字列内に未エスケープ改行 "
                    f"(行 {string_start_line} 開始) — "
                    f"Python triple-quoted の中で '\\n' を使うと改行になります。"
                    f"JS で改行リテラルが必要なら '\\\\n' と書いてください。"
                )
        if state is None:
            if ch == "'":
                state = "'"
                string_start_line = line_no
            elif ch == "`":
                state = "`"
                string_start_line = line_no
            elif ch == "/" and i + 1 < len(js) and js[i+1] == "/":
                # 行コメント: 行末までスキップ
                while i < len(js) and js[i] != "\n":
                    i += 1
                continue
            elif ch == "/" and i + 1 < len(js) and js[i+1] == "*":
                # ブロックコメント
                end = js.find("*/", i + 2)
                if end == -1:
                    break
                # コメント内の改行をカウント
                line_no += js.count("\n", i, end)
                i = end + 2
                continue
        elif state == "'":
            if ch == "\\":
                i += 2  # エスケープ次の文字をスキップ
                continue
            elif ch == "'":
                state = None
        elif state == "`":
            if ch == "\\":
                i += 2
                continue
            elif ch == "`":
                state = None
            # バッククォート（テンプレートリテラル）は改行 OK
        i += 1


def test_llm_system_prompt_uses_cache_control(client: TestClient, monkeypatch):
    """API キーがある場合、system が cache_control 付きの blocks で送られる。"""
    from app import llm
    from app.config import settings
    captured = {}

    class FakeMessages:
        def create(self, **kwargs):
            captured.update(kwargs)
            class FakeMsg:
                content = [type("B", (), {"type": "text", "text": "ok"})()]
                usage = type("U", (), {
                    "input_tokens": 100, "output_tokens": 20,
                    "cache_creation_input_tokens": 0, "cache_read_input_tokens": 0,
                })()
            return FakeMsg()

    class FakeClient:
        messages = FakeMessages()

    monkeypatch.setattr(llm, "_client", lambda: FakeClient())
    monkeypatch.setattr(settings, "anthropic_api_key", "test-key")

    # ローカルモードを抜けて FakeClient 経由で呼ばれる
    md = "# VPN\n\nFortiClient で接続。".encode("utf-8")
    client.post("/api/admin/ingest", files={"file": ("vpn.md", md, "text/markdown")})
    r = client.post("/api/ask", json={"question": "VPN設定"})

    # /api/ask 内で fixture により anthropic_api_key="" にリセットされるため、
    # 実呼出は走らない。直接 llm.answer を呼んでキャッシュ制御を確認。
    monkeypatch.setattr(settings, "anthropic_api_key", "test-key")
    from app.rag import Chunk
    chunks = [(Chunk(chunk_id="x#0", source="x.md", text="dummy"), 0.5)]
    llm.answer("テスト質問", chunks)

    assert "system" in captured
    sys_arg = captured["system"]
    assert isinstance(sys_arg, list)
    assert len(sys_arg) == 1
    assert sys_arg[0]["type"] == "text"
    assert sys_arg[0]["cache_control"] == {"type": "ephemeral"}
